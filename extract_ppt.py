from pptx import Presentation
import json

ppt = Presentation("Dataset_project_repository.pptx")

slides_data = []

for slide_num, slide in enumerate(ppt.slides, start=1):

    slide_text = ""

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            slide_text += shape.text + " "

    slides_data.append({
        "slide": slide_num,
        "text": slide_text
    })

with open("slides.json", "w", encoding="utf-8") as f:
    json.dump(slides_data, f, indent=4)

print("Saved", len(slides_data), "slides")
